from __future__ import annotations

from pathlib import Path


def load_context_files(files: list[Path | str]) -> str:
    """Extract text from a list of PDF, DOCX, PPTX, or TXT files and join them."""
    parts = []
    for f in files:
        path = Path(f)
        text = _extract(path)
        parts.append(f"=== {path.name} ===\n{text.strip()}")
    return "\n\n".join(parts)


def _extract(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _pdf(path)
    if suffix == ".docx":
        return _docx(path)
    if suffix == ".pptx":
        return _pptx(path)
    if suffix in (".txt", ".md"):
        return path.read_text(encoding="utf-8")
    raise ValueError(
        f"Unsupported file type '{suffix}'. "
        "Supported: .pdf, .docx, .pptx, .txt, .md"
    )


def _pdf(path: Path) -> str:
    try:
        import pypdf
    except ImportError:
        raise ImportError(
            "PDF support requires pypdf. Install it with: "
            "pip install 'peer-agents[files]'"
        )
    reader = pypdf.PdfReader(path)
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _docx(path: Path) -> str:
    try:
        import docx
    except ImportError:
        raise ImportError(
            "DOCX support requires python-docx. Install it with: "
            "pip install 'peer-agents[files]'"
        )
    doc = docx.Document(path)
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "PPTX support requires python-pptx. Install it with: "
            "pip install 'peer-agents[files]'"
        )
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            para.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            if para.text.strip()
        ]
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)
